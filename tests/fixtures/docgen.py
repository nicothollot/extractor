"""Synthetic document primitives for Phase-2 tests.

Small builders the reader tests and the realistic memo fixtures share:
text PDFs with ruled (table-finder-detectable) tables, scanned-style PDFs
(text rendered to images — no text layer), image-table pages, rotated and
encrypted PDFs, docx/pptx/xlsx documents. Deterministic output only.
"""

from __future__ import annotations

from pathlib import Path

import fitz

LINE_HEIGHT = 14.0
MARGIN = 54.0
PAGE_WIDTH = 612.0


def table_col_widths(rows: list[list[str]], col_width: float = 95.0) -> list[float]:
    """Per-column widths proportional to content length, clamped to the page.
    Long cells ('Senior Secured Term Loan') must never overflow into the next
    column, and the table must never cross the page edge — either way pymupdf
    merges/truncates the extracted cells."""
    n_cols = max(len(row) for row in rows)
    total = min(n_cols * col_width, PAGE_WIDTH - 2 * MARGIN)
    weights = [
        max(
            max((len(str(row[i])) for row in rows if i < len(row) and row[i]), default=4),
            4,
        )
        for i in range(n_cols)
    ]
    scale = total / sum(weights)
    return [weight * scale for weight in weights]


def draw_table(
    page: fitz.Page,
    origin: tuple[float, float],
    rows: list[list[str]],
    col_width: float = 95.0,
    row_height: float = 18.0,
    fontsize: float = 8.0,
) -> float:
    """Draw a fully ruled table (grid lines + cell text) so pymupdf's table
    finder can detect it. Returns the y just below the table."""
    x0, y0 = origin
    n_rows = len(rows)
    edges = [x0]
    for width in table_col_widths(rows, col_width):
        edges.append(edges[-1] + width)
    x1, y1 = edges[-1], y0 + n_rows * row_height
    for i in range(n_rows + 1):
        y = y0 + i * row_height
        page.draw_line(fitz.Point(x0, y), fitz.Point(x1, y), width=0.5)
    for x in edges:
        page.draw_line(fitz.Point(x, y0), fitz.Point(x, y1), width=0.5)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            if cell:
                page.insert_text(
                    fitz.Point(edges[j] + 3, y0 + i * row_height + row_height - 5),
                    str(cell),
                    fontsize=fontsize,
                )
    return y1


def add_text_page(
    doc: fitz.Document,
    lines: list[str],
    tables: list[list[list[str]]] | None = None,
    rotation: int = 0,
    fontsize: float = 9.0,
) -> fitz.Page:
    """Append one A4-ish page: text lines from the top, then ruled tables."""
    page = doc.new_page(width=612, height=792)
    y = MARGIN
    for line in lines:
        page.insert_text(fitz.Point(MARGIN, y), line, fontsize=fontsize)
        y += LINE_HEIGHT
    for rows in tables or []:
        y += LINE_HEIGHT
        y = draw_table(page, (MARGIN, y), rows) + LINE_HEIGHT
    if rotation:
        page.set_rotation(rotation)
    return page


def make_text_pdf(
    path: Path,
    pages: list[list[str]],
    tables_by_page: dict[int, list[list[list[str]]]] | None = None,
    rotations: dict[int, int] | None = None,
) -> None:
    """Text-layer PDF; pages are 1-based in the tables/rotations maps."""
    doc = fitz.open()
    for number, lines in enumerate(pages, start=1):
        add_text_page(
            doc,
            lines,
            tables=(tables_by_page or {}).get(number),
            rotation=(rotations or {}).get(number, 0),
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    doc.close()


def _render_page_to_png(lines: list[str], tables: list[list[list[str]]] | None = None) -> bytes:
    tmp = fitz.open()
    add_text_page(tmp, lines, tables=tables)
    png = tmp[0].get_pixmap(dpi=150).tobytes("png")
    tmp.close()
    return png


def make_scanned_pdf(path: Path, pages: list[list[str]], fontsize: float = 14.0) -> None:
    """Scanned-style PDF: each page is ONE full-page image, no text layer.
    The default font is bigger than text pages — small print through a
    150dpi 'scan' loses inter-word spacing under OCR."""
    doc = fitz.open()
    for lines in pages:
        tmp = fitz.open()
        add_text_page(tmp, lines, fontsize=fontsize)
        png = tmp[0].get_pixmap(dpi=150).tobytes("png")
        tmp.close()
        page = doc.new_page(width=612, height=792)
        page.insert_image(page.rect, stream=png)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    doc.close()


def _render_table_png(table_rows: list[list[str]]) -> bytes:
    """Render ONLY the table region to a PNG (wide crop)."""
    tmp = fitz.open()
    page = tmp.new_page(width=612, height=792)
    y1 = draw_table(page, (MARGIN, MARGIN), table_rows)
    clip = fitz.Rect(MARGIN - 2, MARGIN - 2, MARGIN + sum(table_col_widths(table_rows)) + 2, y1 + 2)
    png = page.get_pixmap(dpi=150, clip=clip).tobytes("png")
    tmp.close()
    return png


def add_image_table_page(
    doc: fitz.Document, lines: list[str], table_rows: list[list[str]]
) -> fitz.Page:
    """Text page whose table is an embedded IMAGE (wide block, lower half) —
    exercises IMAGE_TABLE classification (D1)."""
    page = doc.new_page(width=612, height=792)
    y = MARGIN
    for line in lines:
        page.insert_text(fitz.Point(MARGIN, y), line, fontsize=9.0)
        y += LINE_HEIGHT
    png = _render_table_png(table_rows)
    target = fitz.Rect(MARGIN, 430, 612 - MARGIN, 700)  # wide, ~34% of the page
    page.insert_image(target, stream=png, keep_proportion=False)
    return page


def make_encrypted_pdf(path: Path, text: str = "secret") -> None:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(fitz.Point(72, 72), text)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(
        str(path),
        encryption=fitz.PDF_ENCRYPT_AES_256,
        user_pw="locked",
        owner_pw="locked",
    )
    doc.close()


def make_docx(
    path: Path,
    sections: list[tuple[str | None, list[str], list[list[list[str]]]]],
) -> None:
    """sections = [(heading or None, paragraphs, tables)]."""
    import docx

    document = docx.Document()
    for heading, paragraphs, tables in sections:
        if heading is not None:
            document.add_heading(heading, level=1)
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
        for rows in tables:
            table = document.add_table(rows=len(rows), cols=max(len(r) for r in rows))
            for i, row in enumerate(rows):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = str(cell) if cell is not None else ""
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(path))


def make_pptx(path: Path, slides: list[tuple[str, list[str], list[list[list[str]]]]]) -> None:
    """slides = [(title, bullet lines, tables)]."""
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    blank = deck.slide_layouts[5]  # title only
    for title, lines, tables in slides:
        slide = deck.slides.add_slide(blank)
        slide.shapes.title.text = title
        if lines:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
            frame = box.text_frame
            frame.text = lines[0]
            for line in lines[1:]:
                frame.add_paragraph().text = line
        top = 4.5
        for rows in tables:
            shape = slide.shapes.add_table(
                len(rows), max(len(r) for r in rows), Inches(0.5), Inches(top), Inches(9), Inches(0.3 * len(rows))
            )
            for i, row in enumerate(rows):
                for j, cell in enumerate(row):
                    shape.table.cell(i, j).text = str(cell) if cell is not None else ""
            top += 0.4 * len(rows)
    path.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(path))


def make_xlsx(path: Path, sheets: dict[str, list[list[object]]]) -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    workbook.remove(workbook.active)
    for name, rows in sheets.items():
        sheet = workbook.create_sheet(title=name)
        for row in rows:
            sheet.append(row)
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(str(path))
