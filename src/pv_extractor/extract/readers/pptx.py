"""PptxReader (D1): python-pptx slide text + tables, for portfolio reviews
delivered as decks. One PageContent per slide (unit_label='slide')."""

from __future__ import annotations

import io
from pathlib import Path

from pv_extractor.models import DocFlag, DocumentContent, PageContent, TableData

from .base import DocumentReader


class PptxReader(DocumentReader):
    name = "pptx"

    def summarize(self, path: str | Path, max_pages: int | None = None) -> DocumentContent:
        try:
            data = self._read_bytes(path)
        except OSError as exc:
            return self._empty(path, DocFlag.ACCESS_ERROR, f"{type(exc).__name__}: {exc}")
        if not data:
            return self._empty(path, DocFlag.CORRUPT_FILE, "zero-byte file")
        try:
            from pptx import Presentation

            deck = Presentation(io.BytesIO(data))
        except Exception as exc:
            return self._empty(path, DocFlag.CORRUPT_FILE, f"{type(exc).__name__}: {exc}")

        pages: list[PageContent] = []
        for number, slide in enumerate(deck.slides, start=1):
            lines: list[str] = []
            tables: list[TableData] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    rows = [
                        [cell.text.strip() or None for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if any(any(cell for cell in row) for row in rows):
                        tables.append(TableData(page_number=number, rows=rows, source="pptx"))
                        lines.extend(
                            " ".join(str(cell) for cell in row if cell) for row in rows
                        )
                elif getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs) or para.text
                        if text.strip():
                            lines.append(text)
            text = "\n".join(lines)
            pages.append(
                PageContent(
                    page_number=number,
                    text=text,
                    tables=tables,
                    text_char_count=len(text.strip()),
                    unit_label="slide",
                    unit_name=getattr(slide.shapes, "title", None) and slide.shapes.title.text or None,
                )
            )
        shown = pages if max_pages is None else pages[:max_pages]
        return DocumentContent(
            file_path=str(path), reader=self.name, page_count=len(pages), pages=shown
        )

    def extract_tables(self, path: str | Path, page_numbers: list[int]) -> dict[int, list[TableData]]:
        content = self.summarize(path)
        wanted = set(page_numbers)
        return {
            page.page_number: page.tables
            for page in content.pages
            if page.page_number in wanted and page.tables
        }
